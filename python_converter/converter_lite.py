"""
Lightweight Document Converter - No heavy ML dependencies
Handles: PDF ↔ Word, PDF ↔ Excel, HTML → PDF, Office → PDF
"""

import sys
import os

def pdf_to_word(input_path, output_path):
    """Convert PDF to Word DOCX with enhanced formatting preservation"""
    from pdf2docx import Converter
    
    cv = Converter(input_path)
    # Enhanced conversion with better layout, formatting, and table detection
    cv.convert(output_path, 
               start=0, 
               end=None,
               pages=None,
               multi_processing=False,
               cpu_count=1)
    cv.close()
    print(f"SUCCESS: Converted PDF to Word: {output_path}")

def word_to_pdf(input_path, output_path):
    """Convert Word DOCX to PDF with formatting preservation"""
    from docx import Document
    from reportlab.lib.pagesizes import letter
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.lib import colors
    
    doc = Document(input_path)
    pdf = SimpleDocTemplate(output_path, pagesize=letter,
                           topMargin=0.75*inch, bottomMargin=0.75*inch,
                           leftMargin=0.75*inch, rightMargin=0.75*inch)
    styles = getSampleStyleSheet()
    story = []
    
    # Add custom styles
    styles.add(ParagraphStyle(name='CustomBold', parent=styles['Normal'], fontName='Helvetica-Bold'))
    styles.add(ParagraphStyle(name='CustomItalic', parent=styles['Normal'], fontName='Helvetica-Oblique'))
    
    # Process paragraphs
    for para in doc.paragraphs:
        if para.text.strip():
            formatted_text = ""
            for run in para.runs:
                text = run.text
                if run.bold:
                    text = f"<b>{text}</b>"
                if run.italic:
                    text = f"<i>{text}</i>"
                if run.underline:
                    text = f"<u>{text}</u>"
                formatted_text += text
            
            p = Paragraph(formatted_text or para.text, styles['Normal'])
            story.append(p)
            story.append(Spacer(1, 0.2*inch))
    
    # Process tables
    for table in doc.tables:
        table_data = [[cell.text for cell in row.cells] for row in table.rows]
        if table_data:
            t = Table(table_data)
            t.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('GRID', (0, 0), (-1, -1), 1, colors.black),
                ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
            ]))
            story.append(t)
            story.append(Spacer(1, 0.3*inch))
    
    pdf.build(story)
    print(f"SUCCESS: Converted Word to PDF: {output_path}")

def pdf_to_excel(input_path, output_path):
    """Extract tables from PDF to Excel with formatting"""
    import pdfplumber
    import openpyxl
    from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
    
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Tables"
    
    current_row = 1
    border = Border(
        left=Side(style='thin'),
        right=Side(style='thin'),
        top=Side(style='thin'),
        bottom=Side(style='thin')
    )
    
    with pdfplumber.open(input_path) as pdf:
        for page_num, page in enumerate(pdf.pages, 1):
            tables = page.extract_tables()
            
            if tables:
                for table in tables:
                    # Page label
                    cell = ws.cell(row=current_row, column=1, value=f"Page {page_num}")
                    cell.font = Font(bold=True, size=12, color='366092')
                    current_row += 1
                    
                    # Table header
                    if table:
                        for col_num, cell_value in enumerate(table[0], 1):
                            cell = ws.cell(row=current_row, column=col_num, value=cell_value)
                            cell.fill = PatternFill(start_color='366092', end_color='366092', fill_type='solid')
                            cell.font = Font(bold=True, color='FFFFFF')
                            cell.alignment = Alignment(horizontal='center', vertical='center')
                            cell.border = border
                        current_row += 1
                        
                        # Table data
                        for row_data in table[1:]:
                            for col_num, cell_value in enumerate(row_data, 1):
                                cell = ws.cell(row=current_row, column=col_num, value=cell_value)
                                cell.border = border
                                cell.alignment = Alignment(horizontal='left', vertical='center')
                            current_row += 1
                    
                    current_row += 1  # Add space between tables
    
    # Auto-adjust column widths
    for column in ws.columns:
        max_length = 0
        column_letter = column[0].column_letter
        for cell in column:
            try:
                if cell.value and len(str(cell.value)) > max_length:
                    max_length = len(str(cell.value))
            except:
                pass
        adjusted_width = min(max_length + 2, 50)
        ws.column_dimensions[column_letter].width = adjusted_width
    
    wb.save(output_path)
    print(f"SUCCESS: Extracted tables to Excel: {output_path}")

def excel_to_pdf(input_path, output_path):
    """Convert Excel XLSX to PDF with enhanced formatting"""
    from openpyxl import load_workbook
    from reportlab.lib.pagesizes import letter, landscape
    from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Spacer, PageBreak
    from reportlab.lib import colors
    from reportlab.lib.units import inch
    
    wb = load_workbook(input_path)
    pagesize = landscape(letter) if wb.active.max_column > 8 else letter
    pdf = SimpleDocTemplate(output_path, pagesize=pagesize,
                           topMargin=0.5*inch, bottomMargin=0.5*inch)
    
    story = []
    
    for sheet_idx, ws in enumerate(wb.worksheets):
        if sheet_idx > 0:
            story.append(PageBreak())
        
        # Get data
        data = []
        for row in ws.iter_rows(values_only=True):
            data.append([str(cell) if cell is not None else '' for cell in row])
            if len(data) >= 100:  # Limit rows
                break
        
        if data:
            t = Table(data)
            t.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#366092')),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, -1), 9),
                ('ALIGN', (0, 0), (-1, 0), 'CENTER'),
                ('ALIGN', (0, 1), (-1, -1), 'LEFT'),
                ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
                ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                ('TOPPADDING', (0, 0), (-1, -1), 6),
                ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
            ]))
            
            # Alternate row colors
            for row_idx in range(1, len(data), 2):
                t.setStyle(TableStyle([
                    ('BACKGROUND', (0, row_idx), (-1, row_idx), colors.HexColor('#F5F5F5'))
                ]))
            
            story.append(t)
            story.append(Spacer(1, 0.3*inch))
    
    pdf.build(story)
    print(f"SUCCESS: Converted Excel to PDF: {output_path}")

def html_to_pdf(input_path, output_path):
    """Convert HTML to PDF with proper table, header, and font support"""
    from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_JUSTIFY
    from html.parser import HTMLParser
    from html import unescape
    import re
    
    class EnhancedHTMLParser(HTMLParser):
        def __init__(self):
            super().__init__()
            self.elements = []
            self.current_text = []
            self.in_bold = False
            self.in_italic = False
            self.in_heading = 0
            self.in_table = False
            self.in_table_row = False
            self.in_table_header = False
            self.current_table = []
            self.current_row = []
            self.current_cell = []
            self.font_stack = []
            
        def handle_starttag(self, tag, attrs):
            attrs_dict = dict(attrs)
            
            if tag == 'table':
                self.in_table = True
                self.current_table = []
            elif tag == 'tr':
                self.in_table_row = True
                self.current_row = []
            elif tag in ['th', 'td']:
                self.in_table_header = (tag == 'th')
                self.current_cell = []
            elif tag in ['b', 'strong']:
                self.in_bold = True
            elif tag in ['i', 'em']:
                self.in_italic = True
            elif tag in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
                self.in_heading = int(tag[1])
            elif tag == 'p':
                if self.current_text:
                    self.flush_text()
            elif tag == 'br':
                self.current_text.append('<br/>')
                
        def handle_endtag(self, tag):
            if tag == 'table':
                if self.current_table:
                    self.elements.append(('table', self.current_table))
                self.in_table = False
            elif tag == 'tr':
                if self.current_row:
                    self.current_table.append(self.current_row)
                self.in_table_row = False
            elif tag in ['th', 'td']:
                cell_text = ' '.join(self.current_cell).strip()
                self.current_row.append({
                    'text': cell_text,
                    'is_header': self.in_table_header
                })
                self.current_cell = []
                self.in_table_header = False
            elif tag in ['b', 'strong']:
                self.in_bold = False
            elif tag in ['i', 'em']:
                self.in_italic = False
            elif tag in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
                self.flush_text()
                self.in_heading = 0
            elif tag in ['p', 'div']:
                self.flush_text()
                
        def handle_data(self, data):
            data = data.strip()
            if not data:
                return
                
            if self.in_table:
                self.current_cell.append(data)
            else:
                # Format text based on current state
                formatted = data
                if self.in_bold:
                    formatted = f'<b>{formatted}</b>'
                if self.in_italic:
                    formatted = f'<i>{formatted}</i>'
                    
                self.current_text.append(formatted)
                
        def flush_text(self):
            if self.current_text:
                text = ' '.join(self.current_text)
                style_type = f'h{self.in_heading}' if self.in_heading else 'normal'
                self.elements.append((style_type, text))
                self.current_text = []
    
    # Read HTML
    with open(input_path, 'r', encoding='utf-8') as f:
        html_content = f.read()
    
    # Parse HTML
    parser = EnhancedHTMLParser()
    parser.feed(html_content)
    parser.flush_text()
    
    # Create PDF
    doc = SimpleDocTemplate(output_path, pagesize=letter,
                           leftMargin=0.75*inch, rightMargin=0.75*inch,
                           topMargin=0.75*inch, bottomMargin=0.75*inch)
    
    styles = getSampleStyleSheet()
    
    # Create custom styles for headings
    for i in range(1, 7):
        size = 18 - (i * 2)
        styles.add(ParagraphStyle(
            name=f'CustomH{i}',
            parent=styles['Heading1'],
            fontSize=size,
            textColor=colors.black,
            spaceAfter=12,
            fontName='Helvetica-Bold'
        ))
    
    story = []
    
    for element_type, content in parser.elements:
        if element_type == 'table':
            # Build table data with proper formatting
            table_data = []
            for row in content:
                row_data = []
                for cell in row:
                    # Create paragraph for cell to support formatting
                    cell_style = styles['Heading4'] if cell['is_header'] else styles['Normal']
                    para = Paragraph(cell['text'], cell_style)
                    row_data.append(para)
                table_data.append(row_data)
            
            # Create table with styling
            if table_data:
                t = Table(table_data)
                t.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, 0), 11),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black),
                    ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                ]))
                story.append(t)
                story.append(Spacer(1, 12))
        elif element_type.startswith('h'):
            # Heading
            level = int(element_type[1])
            para = Paragraph(content, styles[f'CustomH{level}'])
            story.append(para)
            story.append(Spacer(1, 6))
        else:
            # Normal paragraph
            para = Paragraph(content, styles['Normal'])
            story.append(para)
            story.append(Spacer(1, 6))
    
    doc.build(story)
    print(f"SUCCESS: Converted HTML to PDF: {output_path}")

def ppt_to_pdf(input_path, output_path):
    """Convert PowerPoint to PDF"""
    from pptx import Presentation
    from reportlab.lib.pagesizes import letter, landscape
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
    from reportlab.lib.styles import getSampleStyleSheet
    
    prs = Presentation(input_path)
    pdf = SimpleDocTemplate(output_path, pagesize=landscape(letter))
    styles = getSampleStyleSheet()
    story = []
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                p = Paragraph(shape.text, styles['Normal'])
                story.append(p)
                story.append(Spacer(1, 12))
            story.append(PageBreak())
    
    pdf.build(story)
    print(f"SUCCESS: Converted PowerPoint to PDF: {output_path}")

def pdf_to_html(input_path, output_path):
    """Convert PDF to HTML with better formatting, preserving fonts and styles"""
    import pdfplumber
    from html import escape
    
    html_parts = [
        '<!DOCTYPE html>',
        '<html>',
        '<head>',
        '<meta charset="utf-8">',
        '<meta name="viewport" content="width=device-width, initial-scale=1.0">',
        '<style>',
        'body { font-family: Arial, sans-serif; line-height: 1.6; padding: 40px; max-width: 900px; margin: 0 auto; background: #f5f5f5; }',
        '.page { background: white; padding: 50px; margin-bottom: 20px; box-shadow: 0 2px 8px rgba(0,0,0,0.1); page-break-after: always; }',
        '.char { display: inline-block; }',
        '.bold { font-weight: bold; }',
        '.italic { font-style: italic; }',
        '.large { font-size: 1.2em; }',
        '.xlarge { font-size: 1.5em; font-weight: bold; }',
        '.heading { font-size: 1.8em; font-weight: bold; margin: 20px 0 10px 0; }',
        'p { margin: 8px 0; }',
        'table { border-collapse: collapse; width: 100%; margin: 15px 0; }',
        'td, th { border: 1px solid #ddd; padding: 8px; text-align: left; }',
        'th { background-color: #f2f2f2; font-weight: bold; }',
        '@media print { .page { box-shadow: none; margin: 0; } }',
        '</style>',
        '</head>',
        '<body>'
    ]
    
    with pdfplumber.open(input_path) as pdf:
        for page_num, page in enumerate(pdf.pages, 1):
            html_parts.append(f'<div class="page">')
            
            # Extract character-level details for better formatting
            chars = page.chars
            
            if chars:
                # Group characters by line
                lines = {}
                for char in chars:
                    y = round(char['top'])
                    if y not in lines:
                        lines[y] = []
                    lines[y].append(char)
                
                # Process each line
                for y in sorted(lines.keys()):
                    line_chars = sorted(lines[y], key=lambda c: c['x0'])
                    
                    if not line_chars:
                        continue
                    
                    # Detect if line is a heading (large font)
                    avg_size = sum(c.get('size', 12) for c in line_chars) / len(line_chars)
                    is_heading = avg_size > 16
                    
                    line_html = []
                    current_text = []
                    current_style = None
                    
                    for char in line_chars:
                        text = escape(char.get('text', ''))
                        font = char.get('fontname', '').lower()
                        size = char.get('size', 12)
                        
                        # Determine style
                        style_classes = []
                        if 'bold' in font or size > 14:
                            style_classes.append('bold')
                        if 'italic' in font or 'oblique' in font:
                            style_classes.append('italic')
                        if size > 14 and size <= 16:
                            style_classes.append('large')
                        elif size > 16:
                            style_classes.append('xlarge')
                        
                        style = ' '.join(style_classes) if style_classes else None
                        
                        # If style changes, wrap previous text
                        if style != current_style:
                            if current_text:
                                text_content = ''.join(current_text)
                                if current_style:
                                    line_html.append(f'<span class="{current_style}">{text_content}</span>')
                                else:
                                    line_html.append(text_content)
                            current_text = [text]
                            current_style = style
                        else:
                            current_text.append(text)
                    
                    # Add remaining text
                    if current_text:
                        text_content = ''.join(current_text)
                        if current_style:
                            line_html.append(f'<span class="{current_style}">{text_content}</span>')
                        else:
                            line_html.append(text_content)
                    
                    # Wrap in appropriate tag
                    if is_heading:
                        html_parts.append(f'<div class="heading">{"".join(line_html)}</div>')
                    else:
                        html_parts.append(f'<p>{"".join(line_html)}</p>')
            
            # Extract tables
            tables = page.extract_tables()
            if tables:
                for table in tables:
                    html_parts.append('<table>')
                    for i, row in enumerate(table):
                        html_parts.append('<tr>')
                        tag = 'th' if i == 0 else 'td'
                        for cell in row:
                            cell_content = escape(str(cell) if cell else '')
                            html_parts.append(f'<{tag}>{cell_content}</{tag}>')
                        html_parts.append('</tr>')
                    html_parts.append('</table>')
            
            html_parts.append('</div>')
    
    html_parts.extend(['</body>', '</html>'])
    
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write('\n'.join(html_parts))
    
    print(f"SUCCESS: Converted PDF to HTML: {output_path}")

def pdf_to_ppt(input_path, output_path):
    """Convert PDF to PowerPoint"""
    from pptx import Presentation
    from pptx.util import Inches
    import pdfplumber
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    with pdfplumber.open(input_path) as pdf:
        for page in pdf.pages:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            text = page.extract_text()
            
            if text:
                textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(6.5))
                text_frame = textbox.text_frame
                text_frame.text = text
                text_frame.word_wrap = True
    
    prs.save(output_path)
    print(f"SUCCESS: Converted PDF to PowerPoint: {output_path}")

CONVERSIONS = {
    'pdf-to-word': pdf_to_word,
    'word-to-pdf': word_to_pdf,
    'pdf-to-excel': pdf_to_excel,
    'excel-to-pdf': excel_to_pdf,
    'html-to-pdf': html_to_pdf,
    'ppt-to-pdf': ppt_to_pdf,
    'pdf-to-html': pdf_to_html,
    'pdf-to-ppt': pdf_to_ppt,
}

def main():
    if len(sys.argv) != 4:
        print("Usage: converter.exe <conversion_type> <input_file> <output_file>")
        print("\nAvailable conversions:")
        for conv_type in CONVERSIONS.keys():
            print(f"  - {conv_type}")
        sys.exit(1)
    
    conversion_type = sys.argv[1].lower()
    input_file = sys.argv[2]
    output_file = sys.argv[3]
    
    if not os.path.exists(input_file):
        print(f"ERROR: Input file not found: {input_file}")
        sys.exit(1)
    
    if conversion_type not in CONVERSIONS:
        print(f"ERROR: Unknown conversion type: {conversion_type}")
        print(f"Available: {', '.join(CONVERSIONS.keys())}")
        sys.exit(1)
    
    try:
        converter_func = CONVERSIONS[conversion_type]
        converter_func(input_file, output_file)
        sys.exit(0)
    except Exception as e:
        print(f"ERROR: Conversion failed: {str(e)}")
        import traceback
        traceback.print_exc()
        sys.exit(1)

if __name__ == "__main__":
    main()
