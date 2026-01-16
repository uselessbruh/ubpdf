"""
Universal Document Converter
Handles: PDF ↔ Word, PDF ↔ Excel, HTML → PDF, Office → PDF, PDF → Images
Usage: converter.exe <conversion_type> <input_file> <output_file>
"""

import sys
import os
from pathlib import Path
import subprocess

# Path to Ghostscript executable (bundled with app)
def get_ghostscript_path():
    """Get path to bundled Ghostscript executable"""
    if getattr(sys, 'frozen', False):
        # Running as compiled executable
        base_path = os.path.dirname(sys.executable)
    else:
        # Running as script
        base_path = os.path.dirname(os.path.abspath(__file__))
    
    # Look for Ghostscript in parent executables folder
    # Detect architecture
    import platform
    is_64bit = platform.machine().endswith('64')
    gs_folder = 'gs64' if is_64bit else 'gs32'
    gs_exe = 'gswin64c.exe' if is_64bit else 'gswin32c.exe'
    gs_path = os.path.join(os.path.dirname(base_path), 'executables', gs_folder, 'bin', gs_exe)
    if os.path.exists(gs_path):
        return gs_path
    
    # Fallback to same directory as converter
    # Detect architecture
    import platform
    is_64bit = platform.machine().endswith('64')
    gs_folder = 'gs64' if is_64bit else 'gs32'
    gs_exe = 'gswin64c.exe' if is_64bit else 'gswin32c.exe'
    gs_path = os.path.join(base_path, gs_folder, 'bin', gs_exe)
    if os.path.exists(gs_path):
        return gs_path
    
    return None

GHOSTSCRIPT_PATH = get_ghostscript_path()

def pdf_to_word(input_path, output_path):
    """Convert PDF to Word DOCX with enhanced formatting preservation"""
    from pdf2docx import Converter
    
    cv = Converter(input_path)
    # Enhanced conversion with better layout, formatting, image, and table detection
    # Key improvements:
    # - min_section_height: Better section detection for complex layouts
    # - connected_border_tolerance: Better table border detection
    # - min_border_clearance: Improved spacing around borders
    # - float_image_tolerance: Better image positioning
    # - float_layout_tolerance: Improved floating element detection
    cv.convert(output_path, 
               start=0, 
               end=None,
               pages=None,
               multi_processing=False,
               cpu_count=1,
               # Layout settings for better chart/graphic preservation
               min_section_height=20.0,              # Minimum height for section detection (default: 10)
               connected_border_tolerance=0.5,       # Tolerance for connecting borders (default: 0.5)
               min_border_clearance=2.0,             # Minimum clearance around borders (default: 2.0)
               float_image_tolerance=0.1,            # Tolerance for floating images (default: 0.1)
               float_layout_tolerance=0.1,           # Tolerance for floating layout detection (default: 0.1)
               # Image extraction settings
               extract_stream_table=True,            # Extract tables as images if needed (default: False)
               debug=False)
    cv.close()
    print(f"SUCCESS: Converted PDF to Word: {output_path}")

def word_to_pdf(input_path, output_path):
    """Comprehensive Word DOCX to PDF conversion with full formatting preservation"""
    from docx import Document
    from docx.oxml.ns import qn
    from docx.table import Table as DocxTable
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from reportlab.lib.pagesizes import letter
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image as RLImage, PageBreak, ListFlowable, ListItem
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_RIGHT, TA_JUSTIFY
    from reportlab.lib.units import inch
    from reportlab.lib import colors
    from reportlab.pdfgen import canvas
    from PIL import Image
    import re
    import io
    import os
    import tempfile
    
    def rgb_to_hex(rgb_str):
        """Convert Word RGB color to hex color"""
        if not rgb_str or rgb_str == 'auto':
            return None
        try:
            if len(rgb_str) == 6:
                return colors.HexColor(f'#{rgb_str}')
            return None
        except:
            return None
    
    def get_alignment_style(alignment):
        """Convert Word alignment to ReportLab alignment"""
        if alignment == WD_ALIGN_PARAGRAPH.CENTER:
            return TA_CENTER
        elif alignment == WD_ALIGN_PARAGRAPH.RIGHT:
            return TA_RIGHT
        elif alignment == WD_ALIGN_PARAGRAPH.JUSTIFY:
            return TA_JUSTIFY
        return TA_LEFT
    
    doc = Document(input_path)
    pdf = SimpleDocTemplate(output_path, pagesize=letter,
                           topMargin=0.75*inch, bottomMargin=0.75*inch,
                           leftMargin=0.75*inch, rightMargin=0.75*inch)
    styles = getSampleStyleSheet()
    story = []
    
    # Add comprehensive custom styles
    for i in range(1, 7):
        font_sizes = {1: 18, 2: 16, 3: 14, 4: 12, 5: 11, 6: 10}
        styles.add(ParagraphStyle(
            name=f'CustomHeading{i}',
            parent=styles['Heading1'],
            fontSize=font_sizes.get(i, 12),
            fontName='Helvetica-Bold',
            spaceAfter=12,
            spaceBefore=6,
            textColor=colors.HexColor('#2E4053')
        ))
    
    # Add alignment variants
    for align_name, align_val in [('Left', TA_LEFT), ('Center', TA_CENTER), ('Right', TA_RIGHT), ('Justify', TA_JUSTIFY)]:
        styles.add(ParagraphStyle(
            name=f'Normal{align_name}',
            parent=styles['Normal'],
            alignment=align_val
        ))
    
    # Create temp directory for image extraction
    temp_dir = tempfile.mkdtemp()
    image_counter = 0
    
    # Helper function to extract and format run text
    def format_run_text(run):
        """Format a single run with comprehensive styling"""
        text = run.text
        if not text:
            return ""
        
        # Escape special HTML characters
        text = text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
        
        # Apply font size
        if run.font.size:
            size = run.font.size.pt
            text = f'<font size="{int(size)}">{text}</font>'
        
        # Apply text color
        if run.font.color and run.font.color.rgb:
            try:
                color_hex = run.font.color.rgb
                if color_hex:
                    text = f'<font color="#{color_hex}">{text}</font>'
            except:
                pass
        
        # Apply highlight/background color
        if run.font.highlight_color:
            try:
                highlight_colors = {
                    1: '#FFFF00',  # Yellow
                    2: '#00FF00',  # Bright Green
                    3: '#00FFFF',  # Cyan
                    4: '#FF00FF',  # Magenta
                    5: '#0000FF',  # Blue
                    6: '#FF0000',  # Red
                    7: '#000080',  # Dark Blue
                    8: '#008080',  # Teal
                    9: '#008000',  # Green
                    10: '#800080', # Purple
                    11: '#800000', # Dark Red
                    12: '#808000', # Olive
                    13: '#808080', # Gray
                    14: '#C0C0C0', # Light Gray
                    15: '#000000', # Black
                }
                if run.font.highlight_color in highlight_colors:
                    bg_color = highlight_colors[run.font.highlight_color]
                    text = f'<font backColor="{bg_color}">{text}</font>'
            except:
                pass
        
        # Apply superscript and subscript
        if run.font.superscript:
            text = f'<super>{text}</super>'
        elif run.font.subscript:
            text = f'<sub>{text}</sub>'
        
        # Apply bold
        if run.bold:
            text = f'<b>{text}</b>'
        
        # Apply italic
        if run.italic:
            text = f'<i>{text}</i>'
        
        # Apply underline
        if run.underline:
            text = f'<u>{text}</u>'
        
        # Apply strikethrough
        if run.font.strike:
            text = f'<strike>{text}</strike>'
        
        return text
    
    # Process document elements in order
    for element in doc.element.body:
        # Handle paragraphs
        if element.tag.endswith('p'):
            para = None
            for p in doc.paragraphs:
                if p._element == element:
                    para = p
                    break
            
            if para:
                # Check for images in paragraph
                has_image = False
                for run in para.runs:
                    if 'graphic' in run._element.xml:
                        has_image = True
                        try:
                            for rel in run.part.rels.values():
                                if "image" in rel.target_ref:
                                    image_data = rel.target_part.blob
                                    image_counter += 1
                                    
                                    img_path = os.path.join(temp_dir, f"img_{image_counter}.png")
                                    with open(img_path, 'wb') as img_file:
                                        img_file.write(image_data)
                                    
                                    try:
                                        img = RLImage(img_path)
                                        max_width = 6.5 * inch
                                        max_height = 8 * inch
                                        
                                        if img.drawWidth > max_width:
                                            ratio = max_width / img.drawWidth
                                            img.drawWidth = max_width
                                            img.drawHeight = img.drawHeight * ratio
                                        
                                        if img.drawHeight > max_height:
                                            ratio = max_height / img.drawHeight
                                            img.drawHeight = max_height
                                            img.drawWidth = img.drawWidth * ratio
                                        
                                        # Align image based on paragraph alignment
                                        if para.alignment == WD_ALIGN_PARAGRAPH.CENTER:
                                            img.hAlign = 'CENTER'
                                        elif para.alignment == WD_ALIGN_PARAGRAPH.RIGHT:
                                            img.hAlign = 'RIGHT'
                                        else:
                                            img.hAlign = 'LEFT'
                                        
                                        story.append(img)
                                        story.append(Spacer(1, 0.2*inch))
                                    except Exception as e:
                                        pass
                        except Exception as e:
                            pass
                
                # Process text content
                if para.text.strip() and not has_image:
                    # Determine style
                    style_name = 'Normal'
                    base_font_size = 12
                    
                    # Check for heading styles
                    for i in range(1, 7):
                        if para.style.name.startswith(f'Heading {i}'):
                            style_name = f'CustomHeading{i}'
                            break
                    
                    # Get alignment
                    alignment = get_alignment_style(para.alignment) if para.alignment else TA_LEFT
                    
                    # Determine actual style to use
                    if style_name == 'Normal':
                        align_names = {TA_LEFT: 'Left', TA_CENTER: 'Center', TA_RIGHT: 'Right', TA_JUSTIFY: 'Justify'}
                        style_name = f'Normal{align_names.get(alignment, "Left")}'
                    
                    # Get font size from first run
                    if para.runs and para.runs[0].font.size:
                        base_font_size = para.runs[0].font.size.pt
                    
                    # Build formatted text
                    formatted_text = ""
                    for run in para.runs:
                        formatted_text += format_run_text(run)
                    
                    # Create paragraph
                    current_style = styles[style_name]
                    if 'Normal' in style_name:
                        current_style = ParagraphStyle(
                            name=f'{style_name}_temp',
                            parent=current_style,
                            fontSize=base_font_size,
                            leading=base_font_size * 1.2
                        )
                    
                    p = Paragraph(formatted_text or para.text, current_style)
                    story.append(p)
                    story.append(Spacer(1, 0.12*inch))
        
        # Handle tables
        elif element.tag.endswith('tbl'):
            table = None
            for t in doc.tables:
                if t._element == element:
                    table = t
                    break
            
            if table:
                table_data = []
                col_widths = []
                cell_styles = []
                
                for row_idx, row in enumerate(table.rows):
                    row_data = []
                    row_cell_styles = []
                    
                    for col_idx, cell in enumerate(row.cells):
                        # Get cell text with formatting
                        cell_paragraphs = []
                        for para in cell.paragraphs:
                            if para.text.strip():
                                formatted = ""
                                for run in para.runs:
                                    formatted += format_run_text(run)
                                cell_paragraphs.append(formatted)
                        
                        cell_content = '<br/>'.join(cell_paragraphs) if cell_paragraphs else ''
                        
                        # Use plain text for table cells to avoid sizing issues
                        if cell_content:
                            # Strip HTML tags for plain text in tables
                            import re
                            plain_text = re.sub(r'<[^>]+>', '', cell_content)
                            plain_text = plain_text.replace('&amp;', '&').replace('&lt;', '<').replace('&gt;', '>')
                            row_data.append(plain_text)
                        else:
                            row_data.append('')
                        
                        # Track column widths
                        text_len = len(cell.text.strip())
                        if len(col_widths) <= col_idx:
                            col_widths.append(text_len)
                        else:
                            col_widths[col_idx] = max(col_widths[col_idx], text_len)
                        
                        # Get cell background color
                        try:
                            cell_color = None
                            shading = cell._element.xpath('.//w:shd')
                            if shading:
                                fill = shading[0].get(qn('w:fill'))
                                if fill and fill != 'auto':
                                    cell_color = colors.HexColor(f'#{fill}')
                            row_cell_styles.append(cell_color)
                        except:
                            row_cell_styles.append(None)
                    
                    table_data.append(row_data)
                    cell_styles.append(row_cell_styles)
                
                if table_data and len(table_data) > 0:
                    # Calculate column widths
                    total_width = 6.5 * inch
                    total_chars = sum(col_widths) or 1
                    calculated_widths = [(w / total_chars) * total_width for w in col_widths]
                    
                    # Create table
                    t = Table(table_data, colWidths=calculated_widths, repeatRows=1)
                    
                    # Build table style
                    table_style = [
                        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                        ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                        ('FONTSIZE', (0, 0), (-1, -1), 9),
                        ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
                        ('TOPPADDING', (0, 0), (-1, -1), 6),
                        ('LEFTPADDING', (0, 0), (-1, -1), 5),
                        ('RIGHTPADDING', (0, 0), (-1, -1), 5),
                        ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
                    ]
                    
                    # Apply cell background colors
                    for row_idx, row_styles in enumerate(cell_styles):
                        for col_idx, cell_color in enumerate(row_styles):
                            if cell_color:
                                table_style.append(('BACKGROUND', (col_idx, row_idx), (col_idx, row_idx), cell_color))
                    
                    # Default header styling if no custom colors
                    if not cell_styles[0][0]:
                        table_style.append(('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#366092')))
                        table_style.append(('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke))
                    
                    # Alternating row colors for non-header rows
                    for row_idx in range(1, len(table_data)):
                        if not any(cell_styles[row_idx]):
                            bg_color = colors.white if row_idx % 2 == 1 else colors.HexColor('#f9f9f9')
                            table_style.append(('BACKGROUND', (0, row_idx), (-1, row_idx), bg_color))
                    
                    t.setStyle(TableStyle(table_style))
                    story.append(t)
                    story.append(Spacer(1, 0.25*inch))
    
    # Build PDF
    try:
        pdf.build(story)
        print(f"SUCCESS: Converted Word to PDF: {output_path}")
    except Exception as e:
        print(f"ERROR: Failed to build PDF: {str(e)}")
        raise
    finally:
        # Clean up temp directory
        try:
            import shutil
            shutil.rmtree(temp_dir, ignore_errors=True)
        except:
            pass
    
    # Cleanup temp directory
    try:
        import shutil
        shutil.rmtree(temp_dir)
    except:
        pass
    
    print(f"SUCCESS: Converted Word to PDF: {output_path}")
    if image_counter > 0:
        print(f"  Included {image_counter} image(s)")

def pdf_to_excel(input_path, output_path):
    """Extract tables from PDF to Excel - one table per worksheet"""
    import pdfplumber
    import openpyxl
    from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
    
    wb = openpyxl.Workbook()
    wb.remove(wb.active)  # Remove default sheet
    
    border = Border(
        left=Side(style='thin'),
        right=Side(style='thin'),
        top=Side(style='thin'),
        bottom=Side(style='thin')
    )
    
    table_count = 0
    
    with pdfplumber.open(input_path) as pdf:
        for page_num, page in enumerate(pdf.pages, 1):
            tables = page.extract_tables()
            
            if tables:
                for table_idx, table in enumerate(tables, 1):
                    table_count += 1
                    
                    # Create new worksheet for each table
                    # Sheet names: "Table_1", "Table_2", etc. or "Page1_Table1" format
                    sheet_name = f"Page{page_num}_T{table_idx}" if len(tables) > 1 else f"Page{page_num}"
                    # Excel sheet names limited to 31 chars
                    sheet_name = sheet_name[:31]
                    
                    ws = wb.create_sheet(title=sheet_name)
                    current_row = 1
                    
                    # Add page info at top
                    info_cell = ws.cell(row=current_row, column=1, value=f"Source: Page {page_num}, Table {table_idx}")
                    info_cell.font = Font(italic=True, size=10, color='666666')
                    current_row += 1
                    current_row += 1  # Add spacing
                    
                    # Table header
                    if table:
                        for col_num, cell_value in enumerate(table[0], 1):
                            cell = ws.cell(row=current_row, column=col_num, value=cell_value)
                            cell.fill = PatternFill(start_color='366092', end_color='366092', fill_type='solid')
                            cell.font = Font(bold=True, color='FFFFFF', size=11)
                            cell.alignment = Alignment(horizontal='center', vertical='center', wrap_text=True)
                            cell.border = border
                        current_row += 1
                        
                        # Table data
                        for row_data in table[1:]:
                            for col_num, cell_value in enumerate(row_data, 1):
                                cell = ws.cell(row=current_row, column=col_num, value=cell_value)
                                cell.border = border
                                cell.alignment = Alignment(horizontal='left', vertical='center', wrap_text=True)
                            current_row += 1
                    
                    # Auto-adjust column widths
                    for column in ws.columns:
                        max_length = 0
                        column_letter = column[0].column_letter
                        for cell in column:
                            try:
                                if cell.value and len(str(cell.value)) > max_length:
                                    max_length = len(str(cell.value))
                            except:
                                pass
                        adjusted_width = min(max_length + 2, 60)
                        ws.column_dimensions[column_letter].width = adjusted_width
    
    # If no tables found, create a default sheet with message
    if table_count == 0:
        ws = wb.create_sheet(title="No Tables Found")
        ws.cell(row=1, column=1, value="No tables were detected in the PDF file.")
        ws.cell(row=1, column=1).font = Font(italic=True, color='FF0000')
    
    wb.save(output_path)
    print(f"SUCCESS: Extracted {table_count} table(s) to Excel: {output_path}")

def excel_to_pdf(input_path, output_path):
    """Comprehensive Excel XLSX to PDF conversion with intelligent table handling"""
    from openpyxl import load_workbook
    from openpyxl.styles import Font as XLFont, Alignment as XLAlignment
    from openpyxl.utils import get_column_letter
    from reportlab.lib.pagesizes import letter, landscape, A4, A3
    from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, PageBreak
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib import colors
    from reportlab.lib.units import inch
    
    def rgb_to_color(rgb_value):
        """Convert Excel RGB to ReportLab color"""
        if not rgb_value or rgb_value == 'FFFFFFFF' or rgb_value == '00000000':
            return None
        try:
            # Extract RGB from ARGB format (Excel uses ARGB)
            if len(rgb_value) == 8:
                rgb_value = rgb_value[2:]  # Remove alpha channel
            return colors.HexColor(f'#{rgb_value}')
        except:
            return None
    
    def get_alignment(excel_alignment):
        """Convert Excel alignment to ReportLab alignment"""
        if not excel_alignment:
            return 'LEFT'
        alignment_map = {
            'center': 'CENTER',
            'right': 'RIGHT',
            'left': 'LEFT',
            'justify': 'JUSTIFY'
        }
        return alignment_map.get(excel_alignment, 'LEFT')
    
    wb = load_workbook(input_path, data_only=True)
    
    # Analyze all sheets to determine optimal page size and orientation
    max_cols = 0
    max_content_width = 0
    
    for ws in wb.worksheets:
        if ws.max_column > max_cols:
            max_cols = ws.max_column
        
        # Calculate content width based on column widths
        total_width = 0
        for col_idx in range(1, min(ws.max_column + 1, 50)):
            col_letter = get_column_letter(col_idx)
            col_width = ws.column_dimensions[col_letter].width or 10
            total_width += col_width
        
        if total_width > max_content_width:
            max_content_width = total_width
    
    # Choose page size and orientation based on content
    if max_cols > 15 or max_content_width > 200:
        pagesize = landscape(A3)  # Very wide tables - use A3 landscape
    elif max_cols > 10 or max_content_width > 120:
        pagesize = landscape(letter)  # Wide tables - use landscape
    else:
        pagesize = letter  # Normal tables - use portrait
    
    pdf = SimpleDocTemplate(output_path, pagesize=pagesize,
                           topMargin=0.4*inch, bottomMargin=0.4*inch,
                           leftMargin=0.4*inch, rightMargin=0.4*inch)
    
    styles = getSampleStyleSheet()
    story = []
    
    # Process each worksheet
    for sheet_idx, ws in enumerate(wb.worksheets):
        if sheet_idx > 0:
            story.append(PageBreak())
        
        # Add sheet title
        if len(wb.worksheets) > 1:
            title_style = ParagraphStyle('SheetTitle', parent=styles['Heading1'],
                                        fontSize=14, textColor=colors.HexColor('#2C3E50'),
                                        spaceAfter=12, fontName='Helvetica-Bold')
            story.append(Paragraph(f"<b>{ws.title}</b>", title_style))
            story.append(Spacer(1, 0.15*inch))
        
        # Find actual data range (skip empty rows/columns)
        min_row, max_row, min_col, max_col = None, None, None, None
        
        for row in ws.iter_rows():
            for cell in row:
                if cell.value is not None and str(cell.value).strip():
                    if min_row is None or cell.row < min_row:
                        min_row = cell.row
                    if max_row is None or cell.row > max_row:
                        max_row = cell.row
                    if min_col is None or cell.column < min_col:
                        min_col = cell.column
                    if max_col is None or cell.column > max_col:
                        max_col = cell.column
        
        # Skip empty sheets
        if min_row is None or max_row is None:
            story.append(Paragraph("<i>Empty sheet</i>", styles['Normal']))
            continue
        
        # Get data from Excel with proper range
        data = []
        cell_styles = []
        col_widths_chars = []
        
        for row in ws.iter_rows(min_row=min_row, max_row=min(max_row, min_row + 500), 
                                min_col=min_col, max_col=max_col, values_only=False):
            row_data = []
            row_styles = []
            
            for col_idx, cell in enumerate(row):
                # Get cell value with proper formatting
                value = cell.value if cell.value is not None else ''
                
                # Handle numbers, dates, formulas
                if isinstance(value, (int, float)):
                    if cell.number_format and '%' in cell.number_format:
                        value = f"{value:.2%}"
                    elif cell.number_format and any(x in cell.number_format for x in ['0.00', '#,##0']):
                        value = f"{value:,.2f}" if '.' in str(value) else f"{value:,}"
                    else:
                        value = str(value)
                else:
                    value = str(value)
                
                # Create Paragraph for better text wrapping
                cell_para_style = ParagraphStyle(
                    f'Cell_{sheet_idx}_{len(data)}_{col_idx}',
                    parent=styles['Normal'],
                    fontSize=9,
                    leading=11
                )
                
                # Apply cell-level formatting
                if cell.font:
                    if cell.font.bold:
                        cell_para_style.fontName = 'Helvetica-Bold'
                    if cell.font.size:
                        cell_para_style.fontSize = min(cell.font.size, 12)
                        cell_para_style.leading = cell_para_style.fontSize * 1.2
                    if cell.font.color:
                        text_color = rgb_to_color(str(cell.font.color.rgb) if hasattr(cell.font.color, 'rgb') else None)
                        if text_color:
                            cell_para_style.textColor = text_color
                
                # Set alignment
                if cell.alignment:
                    alignment = get_alignment(cell.alignment.horizontal)
                    if alignment == 'CENTER':
                        cell_para_style.alignment = 1
                    elif alignment == 'RIGHT':
                        cell_para_style.alignment = 2
                    else:
                        cell_para_style.alignment = 0
                
                # Wrap value in Paragraph for proper rendering
                if value.strip():
                    para = Paragraph(value, cell_para_style)
                    row_data.append(para)
                else:
                    row_data.append('')
                
                # Track column width
                if len(col_widths_chars) <= col_idx:
                    col_widths_chars.append(len(str(value)))
                else:
                    col_widths_chars[col_idx] = max(col_widths_chars[col_idx], len(str(value)))
                
                # Store cell styling info for background colors
                bg_color = None
                if cell.fill and cell.fill.start_color:
                    bg_color = rgb_to_color(str(cell.fill.start_color.rgb) if hasattr(cell.fill.start_color, 'rgb') else None)
                
                row_styles.append({
                    'bg_color': bg_color,
                    'bold': cell.font.bold if cell.font else False,
                    'is_merged': False
                })
            
            data.append(row_data)
            cell_styles.append(row_styles)
        
        if not data:
            continue
        
        # Calculate intelligent column widths
        available_width = pagesize[0] - 0.8*inch
        num_cols = len(data[0])
        
        # Calculate proportional widths based on content
        total_chars = sum(col_widths_chars) or 1
        col_widths = []
        
        for char_count in col_widths_chars:
            # Proportional width with min/max constraints
            min_width = 0.5 * inch
            max_width = available_width / 2
            
            proportional_width = (char_count / total_chars) * available_width
            calculated_width = max(min_width, min(proportional_width, max_width))
            col_widths.append(calculated_width)
        
        # Ensure total width fits page
        current_total = sum(col_widths)
        if current_total > available_width:
            scale_factor = available_width / current_total
            col_widths = [w * scale_factor for w in col_widths]
        
        # Create table with calculated widths
        t = Table(data, colWidths=col_widths, repeatRows=1)
        
        # Build comprehensive table style
        table_style_commands = [
            ('VALIGN', (0, 0), (-1, -1), 'TOP'),
            ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
            ('TOPPADDING', (0, 0), (-1, -1), 5),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 5),
            ('LEFTPADDING', (0, 0), (-1, -1), 4),
            ('RIGHTPADDING', (0, 0), (-1, -1), 4),
        ]
        
        # Detect and format header row (first row with bold cells or first row by default)
        first_row_bold = any(style.get('bold', False) for style in cell_styles[0])
        
        if first_row_bold or len(data) > 1:
            # Apply header styling to first row
            header_bg = None
            for style in cell_styles[0]:
                if style.get('bg_color'):
                    header_bg = style['bg_color']
                    break
            
            if not header_bg:
                header_bg = colors.HexColor('#2C3E50')
            
            table_style_commands.extend([
                ('BACKGROUND', (0, 0), (-1, 0), header_bg),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('ALIGN', (0, 0), (-1, 0), 'CENTER'),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 8),
                ('TOPPADDING', (0, 0), (-1, 0), 8),
            ])
        
        # Apply cell-specific background colors
        for row_idx, row_style in enumerate(cell_styles):
            for col_idx, cell_style in enumerate(row_style):
                if cell_style.get('bg_color') and row_idx > 0:  # Skip header row
                    table_style_commands.append(
                        ('BACKGROUND', (col_idx, row_idx), (col_idx, row_idx), cell_style['bg_color'])
                    )
        
        # Apply alternating row colors (only to rows without custom background)
        for row_idx in range(1, len(data)):
            row_has_custom_bg = any(style.get('bg_color') for style in cell_styles[row_idx])
            if not row_has_custom_bg:
                bg_color = colors.white if row_idx % 2 == 1 else colors.HexColor('#F8F9FA')
                table_style_commands.append(
                    ('BACKGROUND', (0, row_idx), (-1, row_idx), bg_color)
                )
        
        t.setStyle(TableStyle(table_style_commands))
        story.append(t)
        story.append(Spacer(1, 0.25*inch))
    
    # Build PDF
    try:
        pdf.build(story)
        print(f"SUCCESS: Converted Excel to PDF: {output_path}")
        print(f"  Page size: {pagesize[0]/inch:.1f}\" x {pagesize[1]/inch:.1f}\" ({'Landscape' if pagesize[0] > pagesize[1] else 'Portrait'})")
    except Exception as e:
        print(f"ERROR: Failed to build PDF: {str(e)}")
        raise

def html_to_pdf(input_path, output_path):
    """Comprehensive HTML to PDF conversion with images, tables, lists, and formatting"""
    from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, Image as RLImage, ListFlowable, ListItem, PageBreak
    from reportlab.lib.pagesizes import letter, landscape
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_RIGHT, TA_JUSTIFY
    from html.parser import HTMLParser
    from html import unescape
    from urllib.parse import urljoin, urlparse
    from PIL import Image
    import re
    import os
    import base64
    import io
    import tempfile
    import requests
    
    class ComprehensiveHTMLParser(HTMLParser):
        def __init__(self, base_path):
            super().__init__()
            self.base_path = base_path
            self.elements = []
            self.current_text = []
            self.format_stack = []  # Stack to track nested formatting
            self.in_table = False
            self.in_table_row = False
            self.in_table_header = False
            self.current_table = []
            self.current_row = []
            self.current_cell = []
            self.current_cell_format = []
            self.in_list = False
            self.list_items = []
            self.list_type = 'ul'  # 'ul' or 'ol'
            self.heading_level = 0
            self.paragraph_attrs = {}
            self.temp_images = []
            
        def handle_starttag(self, tag, attrs):
            attrs_dict = dict(attrs)
            
            # Handle images
            if tag == 'img':
                src = attrs_dict.get('src', '')
                alt = attrs_dict.get('alt', 'Image')
                
                if src:
                    self.flush_text()
                    image_data = self.extract_image(src)
                    if image_data:
                        self.elements.append(('image', {
                            'data': image_data,
                            'alt': alt,
                            'align': attrs_dict.get('align', 'left')
                        }))
            
            # Handle tables
            elif tag == 'table':
                self.flush_text()
                self.in_table = True
                self.current_table = {'rows': [], 'attrs': attrs_dict}
                
            elif tag == 'tr':
                self.in_table_row = True
                self.current_row = []
                
            elif tag in ['th', 'td']:
                self.in_table_header = (tag == 'th')
                self.current_cell = []
                self.current_cell_format = []
                
            # Handle lists
            elif tag in ['ul', 'ol']:
                self.flush_text()
                self.in_list = True
                self.list_type = tag
                self.list_items = []
                
            elif tag == 'li':
                if self.in_list:
                    self.current_text = []
            
            # Handle headings
            elif tag in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
                self.flush_text()
                self.heading_level = int(tag[1])
                
            # Handle paragraphs and divs
            elif tag in ['p', 'div']:
                self.flush_text()
                self.paragraph_attrs = {
                    'align': attrs_dict.get('align', 'left'),
                    'style': attrs_dict.get('style', '')
                }
            
            # Handle line breaks
            elif tag == 'br':
                self.current_text.append('<br/>')
            
            # Handle horizontal rules
            elif tag == 'hr':
                self.flush_text()
                self.elements.append(('hr', {}))
            
            # Handle formatting tags
            elif tag in ['b', 'strong']:
                self.format_stack.append('bold')
            elif tag in ['i', 'em']:
                self.format_stack.append('italic')
            elif tag == 'u':
                self.format_stack.append('underline')
            elif tag == 'strike' or tag == 's':
                self.format_stack.append('strike')
            elif tag == 'code' or tag == 'pre':
                self.format_stack.append('code')
            elif tag == 'a':
                href = attrs_dict.get('href', '')
                self.format_stack.append(('link', href))
            elif tag == 'font':
                color = attrs_dict.get('color', '')
                size = attrs_dict.get('size', '')
                self.format_stack.append(('font', color, size))
            elif tag == 'span':
                style = attrs_dict.get('style', '')
                self.format_stack.append(('span', style))
                
        def handle_endtag(self, tag):
            if tag == 'table':
                if self.current_table['rows']:
                    self.elements.append(('table', self.current_table))
                self.in_table = False
                
            elif tag == 'tr':
                if self.current_row:
                    self.current_table['rows'].append(self.current_row)
                self.in_table_row = False
                
            elif tag in ['th', 'td']:
                cell_content = self.format_text(' '.join(self.current_cell).strip())
                self.current_row.append({
                    'text': cell_content,
                    'is_header': self.in_table_header
                })
                self.current_cell = []
                self.current_cell_format = []
                self.in_table_header = False
            
            elif tag in ['ul', 'ol']:
                if self.list_items:
                    self.elements.append(('list', {
                        'type': self.list_type,
                        'items': self.list_items
                    }))
                self.in_list = False
                
            elif tag == 'li':
                if self.in_list and self.current_text:
                    item_text = self.format_text(' '.join(self.current_text).strip())
                    self.list_items.append(item_text)
                    self.current_text = []
            
            elif tag in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
                self.flush_text()
                self.heading_level = 0
                
            elif tag in ['p', 'div']:
                self.flush_text()
                self.paragraph_attrs = {}
            
            # Handle formatting tag closures
            elif tag in ['b', 'strong']:
                if 'bold' in self.format_stack:
                    self.format_stack.remove('bold')
            elif tag in ['i', 'em']:
                if 'italic' in self.format_stack:
                    self.format_stack.remove('italic')
            elif tag == 'u':
                if 'underline' in self.format_stack:
                    self.format_stack.remove('underline')
            elif tag in ['strike', 's']:
                if 'strike' in self.format_stack:
                    self.format_stack.remove('strike')
            elif tag in ['code', 'pre']:
                if 'code' in self.format_stack:
                    self.format_stack.remove('code')
            elif tag == 'a':
                self.format_stack = [f for f in self.format_stack if not (isinstance(f, tuple) and f[0] == 'link')]
            elif tag == 'font':
                self.format_stack = [f for f in self.format_stack if not (isinstance(f, tuple) and f[0] == 'font')]
            elif tag == 'span':
                self.format_stack = [f for f in self.format_stack if not (isinstance(f, tuple) and f[0] == 'span')]
                
        def handle_data(self, data):
            if not data.strip():
                if data:  # Preserve single spaces
                    if self.in_table:
                        self.current_cell.append(' ')
                    else:
                        self.current_text.append(' ')
                return
                
            if self.in_table and not self.in_list:
                self.current_cell.append(data)
            else:
                self.current_text.append(data)
                
        def format_text(self, text):
            """Apply formatting stack to text"""
            if not text:
                return text
            
            # Escape HTML entities
            text = text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            
            for fmt in self.format_stack:
                if fmt == 'bold':
                    text = f'<b>{text}</b>'
                elif fmt == 'italic':
                    text = f'<i>{text}</i>'
                elif fmt == 'underline':
                    text = f'<u>{text}</u>'
                elif fmt == 'strike':
                    text = f'<strike>{text}</strike>'
                elif fmt == 'code':
                    text = f'<font face="Courier">{text}</font>'
                elif isinstance(fmt, tuple):
                    if fmt[0] == 'link':
                        text = f'<font color="blue"><u>{text}</u></font>'
                    elif fmt[0] == 'font':
                        color, size = fmt[1], fmt[2]
                        if color:
                            text = f'<font color="{color}">{text}</font>'
                        if size:
                            text = f'<font size="{size}">{text}</font>'
            
            return text
        
        def flush_text(self):
            if not self.current_text:
                return
                
            text = ' '.join(self.current_text).strip()
            if not text:
                self.current_text = []
                return
            
            formatted_text = self.format_text(text)
            
            if self.heading_level:
                self.elements.append((f'h{self.heading_level}', formatted_text))
            else:
                align = self.paragraph_attrs.get('align', 'left')
                self.elements.append(('p', {'text': formatted_text, 'align': align}))
            
            self.current_text = []
        
        def extract_image(self, src):
            """Extract image from various sources"""
            try:
                # Handle base64 encoded images
                if src.startswith('data:image'):
                    match = re.match(r'data:image/([^;]+);base64,(.+)', src)
                    if match:
                        img_format, img_data = match.groups()
                        img_bytes = base64.b64decode(img_data)
                        return img_bytes
                
                # Handle URLs
                elif src.startswith('http://') or src.startswith('https://'):
                    response = requests.get(src, timeout=10)
                    if response.status_code == 200:
                        return response.content
                
                # Handle local file paths
                else:
                    img_path = src
                    if not os.path.isabs(src):
                        img_path = os.path.join(os.path.dirname(self.base_path), src)
                    
                    if os.path.exists(img_path):
                        with open(img_path, 'rb') as f:
                            return f.read()
            except Exception as e:
                print(f"  Warning: Could not load image: {src} ({e})")
            
            return None
    
    # Read HTML file
    with open(input_path, 'r', encoding='utf-8', errors='ignore') as f:
        html_content = f.read()
    
    # Parse HTML
    parser = ComprehensiveHTMLParser(input_path)
    parser.feed(html_content)
    parser.flush_text()
    
    # Determine page size based on content
    has_wide_tables = any(
        elem[0] == 'table' and len(elem[1]['rows'][0]) > 8 
        for elem in parser.elements 
        if elem[0] == 'table' and elem[1]['rows']
    )
    pagesize = landscape(letter) if has_wide_tables else letter
    
    # Create PDF
    pdf = SimpleDocTemplate(output_path, pagesize=pagesize,
                           topMargin=0.75*inch, bottomMargin=0.75*inch,
                           leftMargin=0.75*inch, rightMargin=0.75*inch)
    
    styles = getSampleStyleSheet()
    
    # Add custom styles
    for i in range(1, 7):
        font_sizes = {1: 20, 2: 16, 3: 14, 4: 12, 5: 11, 6: 10}
        styles.add(ParagraphStyle(
            name=f'CustomH{i}',
            parent=styles['Heading1'],
            fontSize=font_sizes.get(i, 12),
            fontName='Helvetica-Bold',
            spaceAfter=12,
            spaceBefore=6,
            textColor=colors.HexColor('#2C3E50')
        ))
    
    # Alignment styles
    for align_name, align_val in [('Center', TA_CENTER), ('Right', TA_RIGHT), ('Justify', TA_JUSTIFY)]:
        styles.add(ParagraphStyle(
            name=f'Normal{align_name}',
            parent=styles['Normal'],
            alignment=align_val
        ))
    
    story = []
    temp_dir = tempfile.mkdtemp()
    image_counter = 0
    
    # Build PDF content
    for element_type, content in parser.elements:
        if element_type == 'image':
            # Handle images
            try:
                image_counter += 1
                img_path = os.path.join(temp_dir, f"img_{image_counter}.png")
                
                with open(img_path, 'wb') as f:
                    f.write(content['data'])
                
                img = RLImage(img_path)
                max_width = pagesize[0] - 1.5*inch
                max_height = 6*inch
                
                if img.drawWidth > max_width:
                    ratio = max_width / img.drawWidth
                    img.drawWidth = max_width
                    img.drawHeight = img.drawHeight * ratio
                
                if img.drawHeight > max_height:
                    ratio = max_height / img.drawHeight
                    img.drawHeight = max_height
                    img.drawWidth = img.drawWidth * ratio
                
                align = content.get('align', 'left')
                if align == 'center':
                    img.hAlign = 'CENTER'
                elif align == 'right':
                    img.hAlign = 'RIGHT'
                else:
                    img.hAlign = 'LEFT'
                
                story.append(img)
                story.append(Spacer(1, 0.15*inch))
            except Exception as e:
                print(f"  Warning: Could not add image to PDF: {e}")
        
        elif element_type == 'table':
            # Handle tables
            table_rows = content['rows']
            if not table_rows:
                continue
            
            table_data = []
            for row in table_rows:
                row_data = []
                for cell in row:
                    cell_text = cell['text']
                    if cell_text:
                        para = Paragraph(cell_text, styles['Normal'])
                        row_data.append(para)
                    else:
                        row_data.append('')
                table_data.append(row_data)
            
            if table_data:
                # Calculate column widths
                num_cols = len(table_data[0])
                available_width = pagesize[0] - 1.5*inch
                col_widths = [available_width / num_cols] * num_cols
                
                t = Table(table_data, colWidths=col_widths, repeatRows=1)
                
                table_style = [
                    ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                    ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                    ('FONTSIZE', (0, 0), (-1, -1), 9),
                    ('TOPPADDING', (0, 0), (-1, -1), 6),
                    ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
                    ('LEFTPADDING', (0, 0), (-1, -1), 5),
                    ('RIGHTPADDING', (0, 0), (-1, -1), 5),
                    ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
                ]
                
                # Header styling
                if table_rows and table_rows[0] and table_rows[0][0].get('is_header'):
                    table_style.extend([
                        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#2C3E50')),
                        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                        ('ALIGN', (0, 0), (-1, 0), 'CENTER'),
                    ])
                
                # Alternating rows
                for row_idx in range(1, len(table_data)):
                    bg_color = colors.white if row_idx % 2 == 1 else colors.HexColor('#F8F9FA')
                    table_style.append(('BACKGROUND', (0, row_idx), (-1, row_idx), bg_color))
                
                t.setStyle(TableStyle(table_style))
                story.append(t)
                story.append(Spacer(1, 0.2*inch))
        
        elif element_type == 'list':
            # Handle lists
            list_type = content['type']
            items = content['items']
            
            for idx, item in enumerate(items):
                bullet = '•' if list_type == 'ul' else f'{idx + 1}.'
                para_text = f'{bullet} {item}'
                p = Paragraph(para_text, styles['Normal'])
                story.append(p)
                story.append(Spacer(1, 0.05*inch))
            
            story.append(Spacer(1, 0.1*inch))
        
        elif element_type.startswith('h'):
            # Handle headings
            level = element_type[1]
            style_name = f'CustomH{level}'
            p = Paragraph(content, styles.get(style_name, styles['Heading1']))
            story.append(p)
            story.append(Spacer(1, 0.12*inch))
        
        elif element_type == 'p':
            # Handle paragraphs
            text = content['text']
            align = content.get('align', 'left')
            
            style = styles['Normal']
            if align == 'center':
                style = styles.get('NormalCenter', styles['Normal'])
            elif align == 'right':
                style = styles.get('NormalRight', styles['Normal'])
            elif align == 'justify':
                style = styles.get('NormalJustify', styles['Normal'])
            
            if text.strip():
                p = Paragraph(text, style)
                story.append(p)
                story.append(Spacer(1, 0.1*inch))
        
        elif element_type == 'hr':
            # Handle horizontal rules
            from reportlab.platypus import HRFlowable
            story.append(HRFlowable(width="100%", thickness=1, color=colors.grey))
            story.append(Spacer(1, 0.15*inch))
    
    # Build PDF
    try:
        pdf.build(story)
        print(f"SUCCESS: Converted HTML to PDF: {output_path}")
        if image_counter > 0:
            print(f"  Included {image_counter} image(s)")
    except Exception as e:
        print(f"ERROR: Failed to build PDF: {str(e)}")
        raise
    finally:
        # Clean up temp directory
        try:
            import shutil
            shutil.rmtree(temp_dir, ignore_errors=True)
        except:
            pass

def ppt_to_pdf(input_path, output_path):
    """Comprehensive PowerPoint to PDF conversion with full formatting preservation"""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Pt
    from reportlab.lib.pagesizes import letter, landscape
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak, Table, TableStyle, Image as RLImage
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT, TA_JUSTIFY
    from PIL import Image
    import io
    import os
    import tempfile
    
    def rgb_to_color(rgb):
        """Convert RGB tuple to ReportLab color"""
        if rgb is None:
            return None
        try:
            if hasattr(rgb, 'rgb'):
                # Convert from RGBColor object
                r = (rgb.rgb >> 16) & 0xFF
                g = (rgb.rgb >> 8) & 0xFF
                b = rgb.rgb & 0xFF
                return colors.Color(r/255.0, g/255.0, b/255.0)
            return None
        except:
            return None
    
    def get_text_alignment(alignment):
        """Convert PowerPoint alignment to ReportLab alignment"""
        try:
            from pptx.enum.text import PP_ALIGN
            if alignment == PP_ALIGN.CENTER:
                return TA_CENTER
            elif alignment == PP_ALIGN.RIGHT:
                return TA_RIGHT
            elif alignment == PP_ALIGN.JUSTIFY:
                return TA_JUSTIFY
            return TA_LEFT
        except:
            return TA_LEFT
    
    def format_run_text(run):
        """Format a single text run with comprehensive styling"""
        text = run.text
        if not text:
            return ""
        
        # Escape HTML special characters
        text = text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
        
        # Apply font size
        if run.font.size:
            size = run.font.size.pt
            text = f'<font size="{int(size)}">{text}</font>'
        
        # Apply text color
        if run.font.color and run.font.color.rgb:
            try:
                color_obj = run.font.color.rgb
                r = (color_obj >> 16) & 0xFF
                g = (color_obj >> 8) & 0xFF
                b = color_obj & 0xFF
                text = f'<font color="#{r:02x}{g:02x}{b:02x}">{text}</font>'
            except:
                pass
        
        # Apply superscript and subscript
        if run.font.superscript:
            text = f'<super>{text}</super>'
        elif run.font.subscript:
            text = f'<sub>{text}</sub>'
        
        # Apply bold
        if run.font.bold:
            text = f'<b>{text}</b>'
        
        # Apply italic
        if run.font.italic:
            text = f'<i>{text}</i>'
        
        # Apply underline
        if run.font.underline:
            text = f'<u>{text}</u>'
        
        return text
    
    prs = Presentation(input_path)
    
    # Determine page orientation based on slide size
    slide_width = prs.slide_width.inches if prs.slide_width else 10
    slide_height = prs.slide_height.inches if prs.slide_height else 7.5
    pagesize = landscape(letter) if slide_width > slide_height else letter
    
    pdf = SimpleDocTemplate(output_path, pagesize=pagesize,
                           topMargin=0.5*inch, bottomMargin=0.5*inch,
                           leftMargin=0.75*inch, rightMargin=0.75*inch)
    styles = getSampleStyleSheet()
    story = []
    
    # Create temp directory for image extraction
    temp_dir = tempfile.mkdtemp()
    image_counter = 0
    
    # Custom styles
    title_style = ParagraphStyle('CustomTitle', parent=styles['Heading1'],
                                 fontSize=24, textColor=colors.HexColor('#2C3E50'),
                                 alignment=TA_CENTER, spaceAfter=20,
                                 fontName='Helvetica-Bold', spaceBefore=10)
    
    subtitle_style = ParagraphStyle('CustomSubtitle', parent=styles['Heading2'],
                                    fontSize=16, textColor=colors.HexColor('#34495E'),
                                    alignment=TA_CENTER, spaceAfter=15,
                                    fontName='Helvetica')
    
    heading_style = ParagraphStyle('CustomHeading', parent=styles['Heading2'],
                                   fontSize=18, textColor=colors.HexColor('#2C3E50'),
                                   spaceAfter=12, fontName='Helvetica-Bold')
    
    for slide_num, slide in enumerate(prs.slides, 1):
        if slide_num > 1:
            story.append(PageBreak())
        
        # Add slide number with styling
        slide_header = f'<font size=10 color="#7F8C8D"><b>Slide {slide_num}</b></font>'
        slide_label = Paragraph(slide_header, styles['Normal'])
        story.append(slide_label)
        story.append(Spacer(1, 0.15*inch))
        
        # Track if we've found the title
        title_found = False
        
        # Sort shapes by top position to maintain layout order
        shapes_sorted = sorted(slide.shapes, key=lambda s: s.top if hasattr(s, 'top') else 0)
        
        for shape_idx, shape in enumerate(shapes_sorted):
            # Handle pictures/images
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image_stream = shape.image.blob
                    image_counter += 1
                    img_path = os.path.join(temp_dir, f"slide_{slide_num}_img_{image_counter}.png")
                    
                    with open(img_path, 'wb') as img_file:
                        img_file.write(image_stream)
                    
                    # Add image to PDF with proper sizing
                    img = RLImage(img_path)
                    max_width = pagesize[0] - 1.5*inch
                    max_height = 6*inch
                    
                    if img.drawWidth > max_width:
                        ratio = max_width / img.drawWidth
                        img.drawWidth = max_width
                        img.drawHeight = img.drawHeight * ratio
                    
                    if img.drawHeight > max_height:
                        ratio = max_height / img.drawHeight
                        img.drawHeight = max_height
                        img.drawWidth = img.drawWidth * ratio
                    
                    img.hAlign = 'CENTER'
                    story.append(img)
                    story.append(Spacer(1, 0.15*inch))
                except Exception as e:
                    pass
            
            # Handle text boxes and placeholders
            elif hasattr(shape, "text_frame"):
                text_frame = shape.text_frame
                
                # Check if this is a title placeholder
                is_title = False
                is_subtitle = False
                
                if hasattr(shape, 'placeholder_format'):
                    try:
                        from pptx.enum.shapes import PP_PLACEHOLDER
                        placeholder_type = shape.placeholder_format.type
                        if placeholder_type == PP_PLACEHOLDER.TITLE or placeholder_type == PP_PLACEHOLDER.CENTER_TITLE:
                            is_title = True
                            title_found = True
                        elif placeholder_type == PP_PLACEHOLDER.SUBTITLE:
                            is_subtitle = True
                    except:
                        pass
                
                # If not detected as placeholder but is first shape with short text
                if not title_found and shape_idx == 0 and len(shape.text.strip()) < 100:
                    is_title = True
                    title_found = True
                
                # Process paragraphs
                for para in text_frame.paragraphs:
                    if not para.text.strip():
                        continue
                    
                    # Build formatted text
                    formatted_text = ""
                    for run in para.runs:
                        formatted_text += format_run_text(run)
                    
                    if not formatted_text:
                        formatted_text = para.text
                    
                    # Determine style
                    if is_title:
                        p = Paragraph(formatted_text, title_style)
                    elif is_subtitle:
                        p = Paragraph(formatted_text, subtitle_style)
                    else:
                        # Get alignment
                        alignment = get_text_alignment(para.alignment)
                        
                        # Create custom style with alignment
                        para_style = ParagraphStyle(
                            f'TempPara_{slide_num}_{shape_idx}',
                            parent=styles['Normal'],
                            alignment=alignment,
                            fontSize=12,
                            leading=14.4
                        )
                        
                        # Check bullet level
                        if para.level is not None and para.level >= 0:
                            indent = para.level * 25
                            para_style.leftIndent = indent + 15
                            para_style.bulletIndent = indent
                            para_style.spaceAfter = 6
                            
                            # Add bullet based on level
                            bullet_chars = ['•', '◦', '▪', '▫', '–', '·']
                            bullet = bullet_chars[min(para.level, len(bullet_chars)-1)]
                            formatted_text = f'{bullet} {formatted_text}'
                        
                        # Get font size from first run
                        if para.runs and para.runs[0].font.size:
                            para_style.fontSize = para.runs[0].font.size.pt
                            para_style.leading = para_style.fontSize * 1.2
                        
                        p = Paragraph(formatted_text, para_style)
                    
                    story.append(p)
                    story.append(Spacer(1, 0.08*inch))
            
            # Handle tables
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table = shape.table
                table_data = []
                cell_colors = []
                
                for row in table.rows:
                    row_data = []
                    row_colors = []
                    
                    for cell in row.cells:
                        # Get cell text with formatting
                        cell_paragraphs = []
                        for para in cell.text_frame.paragraphs:
                            if para.text.strip():
                                formatted = ""
                                for run in para.runs:
                                    formatted += format_run_text(run)
                                cell_paragraphs.append(formatted)
                        
                        cell_content = '<br/>'.join(cell_paragraphs) if cell_paragraphs else ''
                        
                        if cell_content:
                            cell_para = Paragraph(cell_content, styles['Normal'])
                            row_data.append(cell_para)
                        else:
                            row_data.append('')
                        
                        # Get cell fill color
                        try:
                            fill = cell.fill
                            if fill.type == 1:  # SOLID
                                cell_color = rgb_to_color(fill.fore_color)
                                row_colors.append(cell_color)
                            else:
                                row_colors.append(None)
                        except:
                            row_colors.append(None)
                    
                    table_data.append(row_data)
                    cell_colors.append(row_colors)
                
                if table_data:
                    # Calculate column widths
                    num_cols = len(table_data[0])
                    available_width = pagesize[0] - 1.5*inch
                    col_widths = [available_width / num_cols] * num_cols
                    
                    t = Table(table_data, colWidths=col_widths, repeatRows=1)
                    
                    # Build table style
                    table_style = [
                        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                        ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                        ('FONTSIZE', (0, 0), (-1, -1), 10),
                        ('BOTTOMPADDING', (0, 0), (-1, -1), 8),
                        ('TOPPADDING', (0, 0), (-1, -1), 8),
                        ('LEFTPADDING', (0, 0), (-1, -1), 6),
                        ('RIGHTPADDING', (0, 0), (-1, -1), 6),
                        ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
                    ]
                    
                    # Apply cell background colors
                    for row_idx, row_colors in enumerate(cell_colors):
                        for col_idx, cell_color in enumerate(row_colors):
                            if cell_color:
                                table_style.append(('BACKGROUND', (col_idx, row_idx), (col_idx, row_idx), cell_color))
                    
                    # Default styling if no custom colors
                    if not any(cell_colors[0]):
                        table_style.append(('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#2C3E50')))
                        table_style.append(('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke))
                        table_style.append(('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'))
                    
                    # Alternating row colors
                    for row_idx in range(1, len(table_data)):
                        if not any(cell_colors[row_idx]):
                            bg_color = colors.white if row_idx % 2 == 1 else colors.HexColor('#ECF0F1')
                            table_style.append(('BACKGROUND', (0, row_idx), (-1, row_idx), bg_color))
                    
                    t.setStyle(TableStyle(table_style))
                    story.append(t)
                    story.append(Spacer(1, 0.2*inch))
        
        story.append(Spacer(1, 0.2*inch))
    
    # Build PDF
    try:
        pdf.build(story)
        print(f"SUCCESS: Converted PowerPoint to PDF: {output_path}")
    except Exception as e:
        print(f"ERROR: Failed to build PDF: {str(e)}")
        raise
    finally:
        # Clean up temp directory
        try:
            import shutil
            shutil.rmtree(temp_dir, ignore_errors=True)
        except:
            pass

def pdf_to_html(input_path, output_path):
    """Convert PDF to responsive HTML with images, proper formatting, and mobile support"""
    import pdfplumber
    import fitz  # PyMuPDF for images and detailed formatting
    import os
    from html import escape
    
    # Create directory for images
    output_dir = os.path.dirname(output_path)
    output_name = os.path.splitext(os.path.basename(output_path))[0]
    images_dir = os.path.join(output_dir, f"{output_name}_images")
    os.makedirs(images_dir, exist_ok=True)
    
    # Start HTML with responsive styling
    html_content = """<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <meta http-equiv="X-UA-Compatible" content="ie=edge">
    <title>PDF Document</title>
    <style>
        * {
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }
        
        body {
            font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, 'Helvetica Neue', Arial, sans-serif;
            line-height: 1.6;
            color: #333;
            background-color: #ffffff;
            padding: 20px;
            max-width: 1200px;
            margin: 0 auto;
        }
        
        /* Responsive typography */
        @media (max-width: 768px) {
            body {
                padding: 15px;
                font-size: 16px;
            }
        }
        
        @media (max-width: 480px) {
            body {
                padding: 10px;
                font-size: 14px;
            }
        }
        
        /* Headings */
        h1, h2, h3, h4, h5, h6 {
            margin-top: 1.5em;
            margin-bottom: 0.5em;
            font-weight: 600;
            line-height: 1.3;
        }
        
        h1 { font-size: 2em; }
        h2 { font-size: 1.5em; }
        h3 { font-size: 1.25em; }
        
        @media (max-width: 768px) {
            h1 { font-size: 1.75em; }
            h2 { font-size: 1.35em; }
            h3 { font-size: 1.15em; }
        }
        
        /* Paragraphs and text */
        p {
            margin: 0.8em 0;
            text-align: justify;
        }
        
        .bold { font-weight: bold; }
        .italic { font-style: italic; }
        .underline { text-decoration: underline; }
        
        /* Lists */
        ul, ol {
            margin: 1em 0;
            padding-left: 2em;
        }
        
        li {
            margin: 0.5em 0;
        }
        
        @media (max-width: 768px) {
            ul, ol {
                padding-left: 1.5em;
            }
        }
        
        /* Tables */
        table {
            width: 100%;
            border-collapse: collapse;
            margin: 1.5em 0;
            overflow-x: auto;
            display: block;
        }
        
        @media (max-width: 768px) {
            table {
                font-size: 0.9em;
            }
        }
        
        table thead {
            background-color: #366092;
            color: white;
        }
        
        table th, table td {
            border: 1px solid #ddd;
            padding: 12px;
            text-align: left;
        }
        
        @media (max-width: 768px) {
            table th, table td {
                padding: 8px;
            }
        }
        
        table tr:nth-child(even) {
            background-color: #f9f9f9;
        }
        
        table tr:hover {
            background-color: #f5f5f5;
        }
        
        /* Images */
        img {
            max-width: 100%;
            height: auto;
            display: block;
            margin: 1.5em auto;
            border-radius: 4px;
        }
        
        /* Spacing and structure */
        .section {
            margin-bottom: 2em;
        }
        
        /* Code and preformatted text */
        pre, code {
            background-color: #f4f4f4;
            border: 1px solid #ddd;
            border-radius: 4px;
            padding: 0.2em 0.4em;
            font-family: 'Courier New', Courier, monospace;
            font-size: 0.9em;
        }
        
        pre {
            padding: 1em;
            overflow-x: auto;
        }
        
        /* Links */
        a {
            color: #366092;
            text-decoration: none;
        }
        
        a:hover {
            text-decoration: underline;
        }
        
        /* Blockquotes */
        blockquote {
            border-left: 4px solid #366092;
            padding-left: 1em;
            margin: 1em 0;
            color: #666;
            font-style: italic;
        }
        
        /* Horizontal rules */
        hr {
            border: none;
            border-top: 2px solid #ddd;
            margin: 2em 0;
        }
    </style>
</head>
<body>
"""
    
    # Open PDF with both libraries
    pdf_plumber = pdfplumber.open(input_path)
    pdf_fitz = fitz.open(input_path)
    
    image_counter = 0
    
    for page_num in range(len(pdf_fitz)):
        page_plumber = pdf_plumber.pages[page_num]
        page_fitz = pdf_fitz[page_num]
        
        # Step 1: Extract tables with their bounding boxes using pdfplumber
        table_settings = {
            "vertical_strategy": "lines",
            "horizontal_strategy": "lines",
            "explicit_vertical_lines": [],
            "explicit_horizontal_lines": [],
            "snap_tolerance": 3,
            "join_tolerance": 3,
            "edge_min_length": 3,
            "min_words_vertical": 3,
            "min_words_horizontal": 1,
            "intersection_tolerance": 3,
        }
        
        tables_data = []
        try:
            tables = page_plumber.find_tables(table_settings)
            for table_obj in tables:
                table_bbox = table_obj.bbox  # (x0, top, x1, bottom)
                extracted_table = table_obj.extract()
                
                # Validate table
                if not extracted_table or len(extracted_table) < 2:
                    continue
                
                # Count filled cells
                filled_count = 0
                total_count = 0
                for row in extracted_table:
                    for cell in row:
                        total_count += 1
                        if cell and str(cell).strip():
                            filled_count += 1
                
                # Skip if less than 30% filled or too few cells
                if total_count < 4 or (filled_count / total_count) < 0.3:
                    continue
                
                tables_data.append({
                    'bbox': table_bbox,
                    'data': extracted_table,
                    'y_position': table_bbox[1]  # top coordinate
                })
        except:
            pass
        
        # Step 2: Extract images with positions
        images_data = []
        try:
            image_list = page_fitz.get_images()
            for img_index, img_info in enumerate(image_list):
                try:
                    xref = img_info[0]
                    base_image = pdf_fitz.extract_image(xref)
                    image_bytes = base_image["image"]
                    image_ext = base_image["ext"]
                    
                    # Save image
                    image_counter += 1
                    image_filename = f"image_{image_counter:03d}.{image_ext}"
                    image_path = os.path.join(images_dir, image_filename)
                    
                    with open(image_path, "wb") as img_file:
                        img_file.write(image_bytes)
                    
                    # Get image position
                    img_rects = page_fitz.get_image_rects(xref)
                    if img_rects:
                        rect = img_rects[0]
                        rel_path = f"{output_name}_images/{image_filename}"
                        images_data.append({
                            'bbox': rect,
                            'path': rel_path,
                            'y_position': rect.y0,
                            'counter': image_counter
                        })
                except:
                    pass
        except:
            pass
        
        # Step 3: Extract text blocks, excluding table areas
        text_blocks_data = []
        text_blocks = page_fitz.get_text("dict")
        
        for block in text_blocks.get("blocks", []):
            if block.get("type") != 0:  # Skip non-text blocks
                continue
            
            block_bbox = block.get("bbox")  # (x0, y0, x1, y1)
            if not block_bbox:
                continue
            
            # Check if this block overlaps with any table
            is_in_table = False
            for table_info in tables_data:
                table_bbox = table_info['bbox']  # (x0, top, x1, bottom)
                # Check overlap
                if (block_bbox[0] < table_bbox[2] and block_bbox[2] > table_bbox[0] and
                    block_bbox[1] < table_bbox[3] and block_bbox[3] > table_bbox[1]):
                    is_in_table = True
                    break
            
            if is_in_table:
                continue  # Skip text that's part of a table
            
            # Process text block
            block_html = ""
            for line in block.get("lines", []):
                line_html = ""
                
                for span in line.get("spans", []):
                    text = span.get("text", "").strip()
                    if not text:
                        continue
                    
                    # Get formatting
                    font = span.get("font", "").lower()
                    size = span.get("size", 12)
                    color = span.get("color", 0)
                    
                    # Detect heading based on size
                    if size > 20:
                        tag = "h1"
                    elif size > 16:
                        tag = "h2"
                    elif size > 14:
                        tag = "h3"
                    else:
                        tag = "span"
                    
                    # Build style
                    style_parts = []
                    if size != 12:
                        style_parts.append(f"font-size: {size}pt")
                    
                    if color != 0:
                        r = (color >> 16) & 255
                        g = (color >> 8) & 255
                        b = color & 255
                        style_parts.append(f"color: rgb({r}, {g}, {b})")
                    
                    style_str = "; ".join(style_parts)
                    style_attr = f' style="{style_str}"' if style_str else ""
                    
                    # Escape HTML
                    text_escaped = escape(text)
                    
                    # Apply formatting
                    if "bold" in font:
                        text_escaped = f"<strong>{text_escaped}</strong>"
                    if "italic" in font or "oblique" in font:
                        text_escaped = f"<em>{text_escaped}</em>"
                    
                    # Check for list indicators
                    if text.startswith(('•', '-', '*', '●', '○', '▪', '◆')):
                        line_html = f"<li>{text_escaped.lstrip('•-*●○▪◆ ')}</li>"
                        break
                    elif len(text) > 0 and text[0].isdigit() and '.' in text[:5]:
                        # Numbered list
                        clean_text = '.'.join(text.split('.')[1:]).strip()
                        line_html = f"<li>{escape(clean_text)}</li>"
                        break
                    else:
                        if tag == "span":
                            line_html += f"<{tag}{style_attr}>{text_escaped}</{tag}> "
                        else:
                            line_html = f"<{tag}{style_attr}>{text_escaped}</{tag}>"
                
                if line_html:
                    block_html += line_html + "\n"
            
            if block_html.strip():
                text_blocks_data.append({
                    'html': block_html,
                    'y_position': block_bbox[1],
                    'bbox': block_bbox
                })
        
        # Step 4: Merge and sort all content by Y position
        all_content = []
        
        # Add text blocks
        for text_block in text_blocks_data:
            all_content.append({
                'type': 'text',
                'y_position': text_block['y_position'],
                'html': text_block['html']
            })
        
        # Add tables
        for table_info in tables_data:
            table_html = "<table>\n<thead>\n<tr>\n"
            
            # Header row
            for cell in table_info['data'][0]:
                cell_text = str(cell).strip() if cell else ''
                table_html += f"<th>{escape(cell_text)}</th>\n"
            table_html += "</tr>\n</thead>\n<tbody>\n"
            
            # Data rows
            for row in table_info['data'][1:]:
                table_html += "<tr>\n"
                for cell in row:
                    cell_text = str(cell).strip() if cell else ''
                    table_html += f"<td>{escape(cell_text)}</td>\n"
                table_html += "</tr>\n"
            
            table_html += "</tbody>\n</table>\n"
            
            all_content.append({
                'type': 'table',
                'y_position': table_info['y_position'],
                'html': table_html
            })
        
        # Add images
        for img_info in images_data:
            img_html = f'<img src="{img_info["path"]}" alt="Image {img_info["counter"]}" loading="lazy">\n'
            all_content.append({
                'type': 'image',
                'y_position': img_info['y_position'],
                'html': img_html
            })
        
        # Sort by Y position (top to bottom)
        all_content.sort(key=lambda x: x['y_position'])
        
        # Step 5: Render content in order
        for item in all_content:
            if item['type'] == 'text':
                # Wrap text in appropriate container
                text_html = item['html']
                if "<li>" in text_html:
                    if text_html.count("<li>") > 1:
                        # Multiple list items
                        first_li = text_html[text_html.find("<li>"):text_html.find("</li>")]
                        if any(c.isdigit() for c in first_li[:10]):
                            html_content += f"<ol>\n{text_html}</ol>\n"
                        else:
                            html_content += f"<ul>\n{text_html}</ul>\n"
                    else:
                        html_content += f"<ul>\n{text_html}</ul>\n"
                elif "<h" in text_html:
                    html_content += text_html
                else:
                    html_content += f"<p>{text_html.strip()}</p>\n"
            else:
                # Table or image - already formatted
                html_content += item['html']
    
    pdf_plumber.close()
    pdf_fitz.close()
    
    # Close HTML
    html_content += """
</body>
</html>"""
    
    # Write HTML file
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(html_content)
    
    print(f"SUCCESS: Converted PDF to HTML: {output_path}")
    if image_counter > 0:
        print(f"  Extracted {image_counter} image(s) to: {images_dir}")

def pdf_to_ppt(input_path, output_path):
    """Convert PDF to PowerPoint preserving formatting, images, colors, and layout"""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    import pdfplumber
    from PIL import Image
    import io
    import fitz  # PyMuPDF for better image and color extraction
    
    prs = Presentation()
    
    # Open PDF with both libraries for comprehensive extraction
    pdf_plumber = pdfplumber.open(input_path)
    pdf_fitz = fitz.open(input_path)
    
    for page_num in range(len(pdf_fitz)):
        page_plumber = pdf_plumber.pages[page_num]
        page_fitz = pdf_fitz[page_num]
        
        # Get page dimensions
        page_width = page_fitz.rect.width
        page_height = page_fitz.rect.height
        
        # Set slide dimensions to match PDF page (convert points to inches)
        prs.slide_width = Inches(page_width / 72)
        prs.slide_height = Inches(page_height / 72)
        
        # Create blank slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Extract and set background color if present
        try:
            # Get page background/canvas color
            pixmap = page_fitz.get_pixmap(alpha=False)
            if pixmap.n >= 3:  # RGB or RGBA
                # Sample background color from corners
                img_data = pixmap.tobytes("ppm")
                img = Image.open(io.BytesIO(img_data))
                # Get color from top-left corner (likely background)
                bg_color = img.getpixel((0, 0))
                if bg_color != (255, 255, 255):  # Not white
                    background = slide.background
                    fill = background.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor(bg_color[0], bg_color[1], bg_color[2])
        except:
            pass  # Keep default white background
        
        # Extract images from PDF and add to slide
        try:
            image_list = page_fitz.get_images()
            for img_index, img_info in enumerate(image_list):
                xref = img_info[0]
                base_image = pdf_fitz.extract_image(xref)
                image_bytes = base_image["image"]
                
                # Get image location from page
                img_instances = page_fitz.get_image_rects(xref)
                for img_rect in img_instances:
                    # Convert coordinates from points to inches
                    left = Inches(img_rect.x0 / 72)
                    top = Inches(img_rect.y0 / 72)
                    width = Inches(img_rect.width / 72)
                    height = Inches(img_rect.height / 72)
                    
                    # Save image temporarily
                    img_stream = io.BytesIO(image_bytes)
                    try:
                        slide.shapes.add_picture(img_stream, left, top, width, height)
                    except:
                        pass  # Skip if image format not supported
        except:
            pass  # Continue if image extraction fails
        
        # Extract text with formatting
        text_instances = page_fitz.get_text("dict")
        
        for block in text_instances.get("blocks", []):
            if block.get("type") == 0:  # Text block
                bbox = block.get("bbox")
                if not bbox:
                    continue
                
                # Convert coordinates from points to inches
                left = Inches(bbox[0] / 72)
                top = Inches(bbox[1] / 72)
                width = Inches((bbox[2] - bbox[0]) / 72)
                height = Inches((bbox[3] - bbox[1]) / 72)
                
                # Create textbox
                try:
                    textbox = slide.shapes.add_textbox(left, top, width, height)
                    text_frame = textbox.text_frame
                    text_frame.word_wrap = True
                    text_frame.clear()  # Clear default paragraph
                    
                    # Process each line in the block
                    for line_idx, line in enumerate(block.get("lines", [])):
                        p = text_frame.add_paragraph() if line_idx > 0 else text_frame.paragraphs[0]
                        
                        # Process each span (text segment with same formatting)
                        for span in line.get("spans", []):
                            text = span.get("text", "")
                            if not text.strip():
                                continue
                            
                            run = p.add_run()
                            run.text = text
                            
                            # Apply font formatting
                            font = span.get("font", "")
                            size = span.get("size", 12)
                            color = span.get("color", 0)  # Integer color value
                            
                            # Set font size (convert to points)
                            run.font.size = Pt(size)
                            
                            # Set font color (convert integer to RGB)
                            if color != 0:
                                r = (color >> 16) & 255
                                g = (color >> 8) & 255
                                b = color & 255
                                run.font.color.rgb = RGBColor(r, g, b)
                            
                            # Detect bold/italic from font name
                            if "bold" in font.lower():
                                run.font.bold = True
                            if "italic" in font.lower() or "oblique" in font.lower():
                                run.font.italic = True
                except:
                    pass  # Skip problematic text blocks
        
        # Extract and add tables with formatting
        try:
            tables = page_plumber.extract_tables()
            if tables:
                for table_data in tables:
                    if len(table_data) > 0 and len(table_data[0]) > 0:
                        rows = min(len(table_data), 30)
                        cols = min(len(table_data[0]), 15)
                        
                        # Calculate table position (try to match PDF layout)
                        left = Inches(0.5)
                        top = Inches(1)
                        width = Inches(min(9, prs.slide_width.inches - 1))
                        height = Inches(min(5, prs.slide_height.inches - 2))
                        
                        shape = slide.shapes.add_table(rows, cols, left, top, width, height)
                        ppt_table = shape.table
                        
                        # Fill table data with formatting
                        for i, row in enumerate(table_data[:rows]):
                            for j, cell_value in enumerate(row[:cols]):
                                cell = ppt_table.cell(i, j)
                                cell.text = str(cell_value) if cell_value else ''
                                
                                # Format header row
                                if i == 0:
                                    cell.fill.solid()
                                    cell.fill.fore_color.rgb = RGBColor(54, 96, 146)
                                    for paragraph in cell.text_frame.paragraphs:
                                        for run in paragraph.runs:
                                            run.font.color.rgb = RGBColor(255, 255, 255)
                                            run.font.bold = True
                                            run.font.size = Pt(11)
                                else:
                                    # Regular cells
                                    for paragraph in cell.text_frame.paragraphs:
                                        for run in paragraph.runs:
                                            run.font.size = Pt(10)
        except:
            pass  # Continue if table extraction fails
    
    pdf_plumber.close()
    pdf_fitz.close()
    
    prs.save(output_path)
    print(f"SUCCESS: Converted PDF to PowerPoint: {output_path}")

def pdf_to_images(input_path, output_folder, format='png', dpi=150, quality=90):
    """
    Convert PDF pages to images using Ghostscript for high quality output
    
    Args:
        input_path: Path to input PDF
        output_folder: Folder to save images (if provided) or single image path
        format: Output format ('png', 'jpeg', 'tiff')
        dpi: Resolution in DPI (default: 150)
        quality: JPEG quality 1-100 (default: 90)
    """
    if not GHOSTSCRIPT_PATH or not os.path.exists(GHOSTSCRIPT_PATH):
        raise Exception("Ghostscript not found. PDF to image conversion requires Ghostscript.")
    
    # Check if output_folder is a directory or single file path
    is_single_file = output_folder.lower().endswith(('.png', '.jpg', '.jpeg', '.tiff', '.tif'))
    
    if is_single_file:
        # Single file output
        output_path = output_folder
        output_folder = os.path.dirname(output_path)
        output_basename = os.path.splitext(os.path.basename(output_path))[0]
    else:
        # Directory output
        os.makedirs(output_folder, exist_ok=True)
        output_basename = os.path.splitext(os.path.basename(input_path))[0]
    
    # Map format to Ghostscript device
    device_map = {
        'png': 'png16m',      # 24-bit RGB PNG
        'jpeg': 'jpeg',       # JPEG
        'jpg': 'jpeg',
        'tiff': 'tiff24nc',   # 24-bit RGB TIFF
        'tif': 'tiff24nc'
    }
    
    device = device_map.get(format.lower(), 'png16m')
    ext = format.lower()
    if ext == 'jpg':
        ext = 'jpeg'
    
    # Output file pattern
    output_pattern = os.path.join(output_folder, f"{output_basename}-%d.{ext}")
    
    # Build Ghostscript arguments
    gs_args = [
        GHOSTSCRIPT_PATH,
        '-dNOPAUSE',
        '-dBATCH',
        '-dSAFER',
        '-sDEVICE=' + device,
        f'-r{dpi}',           # Set resolution
        '-dTextAlphaBits=4',  # Anti-aliasing for text
        '-dGraphicsAlphaBits=4',  # Anti-aliasing for graphics
    ]
    
    # Add JPEG quality if applicable
    if format.lower() in ['jpeg', 'jpg']:
        gs_args.append(f'-dJPEGQ={quality}')
    
    gs_args.extend([
        f'-sOutputFile={output_pattern}',
        input_path
    ])
    
    # Execute Ghostscript
    try:
        result = subprocess.run(gs_args, capture_output=True, text=True, check=True)
        
        # If single file output, rename the first generated file
        if is_single_file:
            first_generated = os.path.join(output_folder, f"{output_basename}-1.{ext}")
            if os.path.exists(first_generated):
                os.rename(first_generated, output_path)
                # Remove any additional pages if they exist
                page_num = 2
                while True:
                    next_page = os.path.join(output_folder, f"{output_basename}-{page_num}.{ext}")
                    if os.path.exists(next_page):
                        os.remove(next_page)
                        page_num += 1
                    else:
                        break
        
        print(f"SUCCESS: Converted PDF to {format.upper()} images: {output_folder if not is_single_file else output_path}")
        return True
    except subprocess.CalledProcessError as e:
        raise Exception(f"Ghostscript conversion failed: {e.stderr}")

def pdf_to_png(input_path, output_path):
    """Convert PDF to PNG using Ghostscript"""
    return pdf_to_images(input_path, output_path, format='png', dpi=150)

def pdf_to_jpeg(input_path, output_path):
    """Convert PDF to JPEG using Ghostscript"""
    return pdf_to_images(input_path, output_path, format='jpeg', dpi=150, quality=90)

def pdf_to_tiff(input_path, output_path):
    """Convert PDF to TIFF using Ghostscript"""
    return pdf_to_images(input_path, output_path, format='tiff', dpi=150)

# Conversion mapping
CONVERSIONS = {
    'pdf-to-word': pdf_to_word,
    'word-to-pdf': word_to_pdf,
    'pdf-to-excel': pdf_to_excel,
    'excel-to-pdf': excel_to_pdf,
    'html-to-pdf': html_to_pdf,
    'ppt-to-pdf': ppt_to_pdf,
    'pdf-to-html': pdf_to_html,
    'pdf-to-ppt': pdf_to_ppt,
    'pdf-to-png': pdf_to_png,
    'pdf-to-jpeg': pdf_to_jpeg,
    'pdf-to-jpg': pdf_to_jpeg,
    'pdf-to-tiff': pdf_to_tiff,
    'pdf-to-images': pdf_to_images,
}

def main():
    if len(sys.argv) != 4:
        print("Usage: converter.exe <conversion_type> <input_file> <output_file>")
        print("\nAvailable conversions:")
        for conv_type in CONVERSIONS.keys():
            print(f"  - {conv_type}")
        sys.exit(1)
    
    conversion_type = sys.argv[1].lower()
    input_file = sys.argv[2]
    output_file = sys.argv[3]
    
    # Validate input file exists
    if not os.path.exists(input_file):
        print(f"ERROR: Input file not found: {input_file}")
        sys.exit(1)
    
    # Get conversion function
    if conversion_type not in CONVERSIONS:
        print(f"ERROR: Unknown conversion type: {conversion_type}")
        print(f"Available: {', '.join(CONVERSIONS.keys())}")
        sys.exit(1)
    
    try:
        # Run conversion
        converter_func = CONVERSIONS[conversion_type]
        converter_func(input_file, output_file)
        sys.exit(0)
    except Exception as e:
        print(f"ERROR: Conversion failed: {str(e)}")
        import traceback
        traceback.print_exc()
        sys.exit(1)

if __name__ == "__main__":
    main()
